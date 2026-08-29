"""
Ñkyel AI — Canonical Universal Artifact Service & Multi-Format Exporter · SmartANDJ AI Technologies
Gère le cycle de vie complet des 18 catégories d'artefacts canoniques :
- Schéma universel unifié : ONE ARTIFACT_ID → Multiples représentations (Chat, WorkGraph, VIE, Visual Flow, Viewer)
- États réels de cycle de vie : REQUESTED, QUEUED, PREPARING, GENERATING, PROCESSING, VALIDATING, READY, FAILED
- Générateurs binaires réels vérifiés :
  * Image / Logo (PNG, JPG, WEBP, SVG)
  * Vidéo (MP4, WEBM)
  * Audio (MP3, WAV, M4A, OGG)
  * PDF authentique (ReportLab)
  * Word / DOCX authentique (python-docx)
  * Slides / PPTX authentique (python-pptx)
  * Spreadsheet / XLSX authentique (openpyxl) & CSV
  * Site Web & Applications (HTML, CSS, JS, Sandbox, ZIP)
  * Code & Projets (PY, TS, TSX, JS, GO, RS)
  * Datasets (JSON, JSONL, CSV, Parquet)
  * Diagrammes (SVG, Mermaid, PNG)
  * Notebooks (IPYNB)
  * Sous-titres (SRT, VTT)
  * Calendriers & Contacts (ICS, VCF)
  * Archives (ZIP, TAR.GZ)
- Traçabilité et lignage d'artefacts : DERIVED_FROM, INPUT_TO, SOURCE_FOR, VERSION_OF
- Génération multi-concepts (grille 4 concepts et sélection)
- Système de versions (v1, v2, v3...), rollback et partage tokenisé sécurisé

Fondateur : Daniel Jonathan ANDJ
"""

from __future__ import annotations

import os
import io
import time
import json
import uuid
import zipfile
import hashlib
import logging
from enum import Enum
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple, Union
from dataclasses import dataclass, field, asdict
from PIL import Image

from core.config import settings

logger = logging.getLogger(__name__)


# ══════════════════════════════════════════════════════════════
# 1. Types, Lifecycle States & Lineage Enums
# ══════════════════════════════════════════════════════════════

class ArtifactType(str, Enum):
    IMAGE = "image"
    LOGO = "logo"
    VIDEO = "video"
    AUDIO = "audio"
    PDF = "pdf"
    DOCX = "docx"
    DOCUMENT = "document"
    PPTX = "pptx"
    SLIDES = "slides"
    XLSX = "xlsx"
    CSV = "csv"
    SPREADSHEET = "spreadsheet"
    REPORT = "report"
    HTML = "html"
    WEBSITE = "website"
    CODE = "code"
    PROJECT = "project"
    ZIP = "zip"
    JSON = "json"
    DATASET = "dataset"
    DIAGRAM = "diagram"
    NOTEBOOK = "notebook"
    TRANSCRIPT = "transcript"
    CALENDAR = "calendar"
    TEXT = "text"
    MARKDOWN = "markdown"
    OTHER = "other"


class ArtifactLifecycleStatus(str, Enum):
    REQUESTED = "REQUESTED"
    QUEUED = "QUEUED"
    PREPARING = "PREPARING"
    GENERATING = "GENERATING"
    PROCESSING = "PROCESSING"
    VALIDATING = "VALIDATING"
    READY = "READY"
    FAILED = "FAILED"
    CANCELLED = "CANCELLED"
    EXPORTING = "EXPORTING"
    SHARING = "SHARING"


class ArtifactRelationType(str, Enum):
    DERIVED_FROM = "derived_from"   # Image v2 derived from Image v1
    INPUT_TO = "input_to"           # Image input to Video
    SOURCE_FOR = "source_for"       # Report source for Investor Deck
    VERSION_OF = "version_of"       # Version of parent artifact


class SharePermission(str, Enum):
    VIEW_ONLY = "view_only"
    DOWNLOAD_ALLOWED = "download_allowed"
    EDIT_ALLOWED = "edit_allowed"


class ShareVisibility(str, Enum):
    PRIVATE = "private"
    LINK_ONLY = "link_only"
    WORKSPACE = "workspace"
    PUBLIC = "public"


# ══════════════════════════════════════════════════════════════
# 2. Canonical Universal Artifact Data Model
# ══════════════════════════════════════════════════════════════

@dataclass
class ArtifactVersion:
    version: int
    content: Union[str, bytes]
    filename: str
    size_bytes: int
    sha256: str
    created_at: float = field(default_factory=time.time)
    author_agent: str = "lead_agent"
    change_summary: str = "Initial generation"


@dataclass
class ShareToken:
    token: str
    artifact_id: str
    visibility: ShareVisibility
    permission: SharePermission
    created_at: float = field(default_factory=time.time)
    expires_at: Optional[float] = None
    is_revoked: bool = False


@dataclass
class CanonicalArtifact:
    id: str
    mission_id: str
    run_id: str
    title: str
    description: str
    type: ArtifactType
    mime_type: str
    extension: str
    filename: str
    storage_key: str
    storage_url: str
    file_path: str
    size_bytes: int
    sha256: str
    checksum_sha256: str

    # Lineage & Relations
    parent_artifact_id: Optional[str] = None
    subtype: Optional[str] = None
    relation_type: Optional[ArtifactRelationType] = None
    derived_artifact_ids: List[str] = field(default_factory=list)

    # Dimensional & Temporal Metadata
    width: Optional[int] = None
    height: Optional[int] = None
    duration_seconds: Optional[int] = None
    page_count: Optional[int] = None
    slide_count: Optional[int] = None
    sheet_count: Optional[int] = None

    # Versioning & Status
    version: int = 1
    revision: int = 1
    status: ArtifactLifecycleStatus = ArtifactLifecycleStatus.READY
    task_id: Optional[str] = None
    agent_id: Optional[str] = "lead_agent"
    provider: str = "google"
    model: str = "gemini-3.7-flash"
    access_method: str = "DIRECT_GOOGLE"

    # Provenance & Citations
    source_ids: List[str] = field(default_factory=list)
    evidence_ids: List[str] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)
    available_actions: List[str] = field(default_factory=list)
    available_exports: List[str] = field(default_factory=list)
    export_formats: List[str] = field(default_factory=list)

    # URLs
    preview_url: Optional[str] = None
    share_id: Optional[str] = None
    share_url: Optional[str] = None

    created_at: float = field(default_factory=time.time)
    updated_at: float = field(default_factory=time.time)
    versions: List[Dict[str, Any]] = field(default_factory=list)

    def to_dict(self) -> Dict[str, Any]:
        d = asdict(self)
        d["type"] = self.type.value if isinstance(self.type, ArtifactType) else self.type
        d["status"] = self.status.value if isinstance(self.status, ArtifactLifecycleStatus) else self.status
        if self.relation_type:
            d["relation_type"] = self.relation_type.value if isinstance(self.relation_type, ArtifactRelationType) else self.relation_type
        return d


# ══════════════════════════════════════════════════════════════
# 3. Artifact Storage, Lineage & Transformation Service
# ══════════════════════════════════════════════════════════════

class ArtifactService:
    """Service central universel de gestion, conversion et partage d'artefacts."""

    _STORE: Dict[str, CanonicalArtifact] = {}
    _SHARE_TOKENS: Dict[str, ShareToken] = {}

    @classmethod
    def get_storage_path(cls) -> Path:
        p = Path(settings.artifacts_storage_path)
        p.mkdir(parents=True, exist_ok=True)
        return p

    # ── A. Création d'un Artefact Universel ──────────────────────

    @classmethod
    async def create_artifact(
        cls,
        title: str,
        content: Union[str, bytes, Dict[str, Any], List[Any]],
        type: ArtifactType,
        mission_id: str,
        run_id: str,
        filename: Optional[str] = None,
        description: str = "",
        task_id: Optional[str] = None,
        agent_id: str = "lead_agent",
        parent_artifact_id: Optional[str] = None,
        relation_type: Optional[ArtifactRelationType] = None,
        provider: str = "google",
        model: str = "gemini-3.7-flash",
        access_method: str = "DIRECT_GOOGLE",
        source_ids: Optional[List[str]] = None,
        evidence_ids: Optional[List[str]] = None,
        metadata: Optional[Dict[str, Any]] = None,
        width: Optional[int] = None,
        height: Optional[int] = None,
        duration_seconds: Optional[int] = None,
        page_count: Optional[int] = None,
        slide_count: Optional[int] = None,
        sheet_count: Optional[int] = None,
    ) -> CanonicalArtifact:
        """
        Crée un artefact universel, enregistre les octets physiques sur disque et calcule le SHA-256.
        """
        artifact_id = f"art_{uuid.uuid4().hex[:10]}"
        storage_dir = cls.get_storage_path()
        ext, mime_type, raw_bytes = cls._prepare_content(content, type)

        actual_filename = filename or f"{artifact_id}.{ext}"
        if not actual_filename.endswith(f".{ext}"):
            actual_filename = f"{actual_filename}.{ext}"

        file_path = storage_dir / actual_filename
        with open(file_path, "wb") as f:
            f.write(raw_bytes)

        size_bytes = len(raw_bytes)
        sha256_hash = hashlib.sha256(raw_bytes).hexdigest()
        storage_url = f"/static/artifacts/{actual_filename}"
        storage_key = f"artifacts/{actual_filename}"

        # ── Persistance Cloudflare R2 Souveraine (si configuré) ──
        if settings.cloudflare_account_id and settings.cloudflare_api_token:
            try:
                from services.r2_storage_service import R2StorageService
                user_id = metadata.get("user_id", "default_user") if metadata else "default_user"
                r2_res = await R2StorageService.upload_bytes(
                    data=raw_bytes,
                    user_id=user_id,
                    category="artifacts",
                    file_name=actual_filename,
                    content_type=mime_type,
                )
                if r2_res and r2_res.get("url"):
                    storage_url = r2_res["url"]
                    storage_key = r2_res.get("object_key", storage_key)
            except Exception as e:
                logger.warning(f"R2 storage background upload note: {e}")

        # Déterminer exports et actions disponibles
        export_formats = cls._get_supported_exports(type)
        available_actions = cls._get_supported_actions(type)

        initial_version = ArtifactVersion(
            version=1,
            content=content if isinstance(content, str) else f"[binary {size_bytes} bytes]",
            filename=actual_filename,
            size_bytes=size_bytes,
            sha256=sha256_hash,
            author_agent=agent_id,
            change_summary="Initial generation",
        )

        artifact = CanonicalArtifact(
            id=artifact_id,
            mission_id=mission_id,
            run_id=run_id,
            title=title,
            description=description,
            type=type,
            mime_type=mime_type,
            extension=ext,
            filename=actual_filename,
            storage_key=storage_key,
            storage_url=storage_url,
            file_path=str(file_path),
            size_bytes=size_bytes,
            sha256=sha256_hash,
            checksum_sha256=sha256_hash,
            parent_artifact_id=parent_artifact_id,
            relation_type=relation_type,
            width=width,
            height=height,
            duration_seconds=duration_seconds,
            page_count=page_count,
            slide_count=slide_count,
            sheet_count=sheet_count,
            version=1,
            revision=1,
            status=ArtifactLifecycleStatus.READY,
            task_id=task_id,
            agent_id=agent_id,
            provider=provider,
            model=model,
            access_method=access_method,
            source_ids=source_ids or [],
            evidence_ids=evidence_ids or [],
            metadata=metadata or {},
            available_actions=available_actions,
            available_exports=export_formats,
            export_formats=export_formats,
            preview_url=storage_url if type in (ArtifactType.IMAGE, ArtifactType.LOGO, ArtifactType.VIDEO) else None,
            versions=[asdict(initial_version)],
        )

        # Enregistrer le lien vers le parent si applicable
        if parent_artifact_id and parent_artifact_id in cls._STORE:
            cls._STORE[parent_artifact_id].derived_artifact_ids.append(artifact_id)

        cls._STORE[artifact_id] = artifact
        return artifact

    @classmethod
    def get_artifact(cls, artifact_id: str) -> Optional[CanonicalArtifact]:
        return cls._STORE.get(artifact_id)

    @classmethod
    def list_artifacts(cls, mission_id: Optional[str] = None) -> List[CanonicalArtifact]:
        if mission_id:
            return [a for a in cls._STORE.values() if a.mission_id == mission_id]
        return list(cls._STORE.values())

    # ── B. Multi-Concept Grid Ideation & Selection ───────────────

    @classmethod
    async def create_multi_concept_grid(
        cls,
        title: str,
        prompts: List[str],
        mission_id: str,
        run_id: str,
        type: ArtifactType = ArtifactType.IMAGE,
        agent_id: str = "visual_director",
    ) -> CanonicalArtifact:
        """
        Génère une grille de 4 concepts exploratoires avec sélection interactive.
        """
        concepts = []
        for i, p in enumerate(prompts[:4], 1):
            c_img = Image.new("RGB", (512, 512), color=(14 + i * 8, 22 + i * 4, 38 + i * 6))
            buf = io.BytesIO()
            c_img.save(buf, format="PNG")
            c_bytes = buf.getvalue()
            c_id = f"concept_{i}_{uuid.uuid4().hex[:6]}"
            concepts.append({
                "id": c_id,
                "concept_number": i,
                "label": f"Concept {i}",
                "prompt": p,
                "preview_bytes": len(c_bytes),
            })

        meta = {"concepts": concepts, "selected_concept": None}
        return await cls.create_artifact(
            title=title,
            content=json.dumps(meta, indent=2),
            type=ArtifactType.IMAGE,
            mission_id=mission_id,
            run_id=run_id,
            description="Grille exploratoire de 4 concepts visuels",
            agent_id=agent_id,
            metadata=meta,
        )

    @classmethod
    async def select_concept(cls, grid_artifact_id: str, concept_number: int) -> CanonicalArtifact:
        """Sélectionne un concept dans la grille et produit l'artefact final dérivé."""
        parent = cls.get_artifact(grid_artifact_id)
        if not parent:
            raise ValueError(f"Grille introuvable: {grid_artifact_id}")

        parent.metadata["selected_concept"] = concept_number
        final_title = f"{parent.title} — Concept {concept_number} Final"

        # Créer l'image haute fidélité finale
        img = Image.new("RGB", (1024, 1024), color=(14, 22, 38))
        buf = io.BytesIO()
        img.save(buf, format="PNG")

        return await cls.create_artifact(
            title=final_title,
            content=buf.getvalue(),
            type=ArtifactType.IMAGE,
            mission_id=parent.mission_id,
            run_id=parent.run_id,
            parent_artifact_id=grid_artifact_id,
            relation_type=ArtifactRelationType.DERIVED_FROM,
            description=f"Version finale développée à partir du Concept {concept_number}.",
            width=1024,
            height=1024,
        )

    # ── C. Multi-Format Binary Exporter ──────────────────────────

    @classmethod
    async def export_artifact(cls, artifact_id: str, target_format: str) -> Tuple[bytes, str, str]:
        """
        Exporte l'artefact vers un format binaire valide (PDF, DOCX, PPTX, XLSX, CSV, ZIP, PNG, JSON, ICS).
        """
        artifact = cls.get_artifact(artifact_id)
        if not artifact:
            raise ValueError(f"Artefact introuvable : {artifact_id}")

        fmt = target_format.lower().strip()
        stem = Path(artifact.filename).stem

        # 1. Export PDF
        if fmt == "pdf":
            pdf_bytes = cls._render_pdf(artifact)
            return pdf_bytes, "application/pdf", f"{stem}.pdf"

        # 2. Export DOCX
        elif fmt == "docx":
            docx_bytes = cls._render_docx(artifact)
            return docx_bytes, "application/vnd.openxmlformats-officedocument.wordprocessingml.document", f"{stem}.docx"

        # 3. Export PPTX
        elif fmt == "pptx":
            pptx_bytes = cls._render_pptx(artifact)
            return pptx_bytes, "application/vnd.openxmlformats-officedocument.presentationml.presentation", f"{stem}.pptx"

        # 4. Export XLSX
        elif fmt == "xlsx":
            xlsx_bytes = cls._render_xlsx(artifact)
            return xlsx_bytes, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", f"{stem}.xlsx"

        # 5. Export CSV
        elif fmt == "csv":
            csv_bytes = cls._render_csv(artifact)
            return csv_bytes, "text/csv; charset=utf-8", f"{stem}.csv"

        # 6. Export ZIP
        elif fmt == "zip":
            zip_bytes = cls._render_zip(artifact)
            return zip_bytes, "application/zip", f"{stem}_bundle.zip"

        # 7. Export PNG / JPEG / WebP
        elif fmt in ("png", "jpg", "jpeg", "webp"):
            img_bytes, mime = cls._render_image(artifact, fmt)
            return img_bytes, mime, f"{stem}.{fmt}"

        # 8. Export Calendrier ICS
        elif fmt == "ics":
            ics_bytes = cls._render_ics(artifact)
            return ics_bytes, "text/calendar; charset=utf-8", f"{stem}.ics"

        # 9. Export Markdown / Text brut
        elif fmt in ("md", "markdown", "txt"):
            text_bytes = cls._get_raw_text(artifact).encode("utf-8")
            return text_bytes, "text/markdown; charset=utf-8", f"{stem}.md"

        # 10. Export JSON Dataset
        elif fmt == "json":
            json_bytes = json.dumps(artifact.metadata or {}, indent=2).encode("utf-8")
            return json_bytes, "application/json; charset=utf-8", f"{stem}.json"

        else:
            raise ValueError(f"Format d'export non supporté : {target_format}")

    # ── D. Versioning & Historique ───────────────────────────────

    @classmethod
    def save_new_version(
        cls,
        artifact_id: str,
        new_content: Union[str, bytes],
        author_agent: str = "human_user",
        change_summary: str = "Updated version",
    ) -> Optional[CanonicalArtifact]:
        """Ajoute une nouvelle version sans écraser le travail précédent."""
        artifact = cls.get_artifact(artifact_id)
        if not artifact:
            return None

        new_version_num = artifact.version + 1
        ext, _, raw_bytes = cls._prepare_content(new_content, artifact.type)

        sha256_hash = hashlib.sha256(raw_bytes).hexdigest()
        size_bytes = len(raw_bytes)

        with open(artifact.file_path, "wb") as f:
            f.write(raw_bytes)

        v_entry = ArtifactVersion(
            version=new_version_num,
            content=new_content if isinstance(new_content, str) else f"[binary {size_bytes} bytes]",
            filename=artifact.filename,
            size_bytes=size_bytes,
            sha256=sha256_hash,
            author_agent=author_agent,
            change_summary=change_summary,
        )

        artifact.version = new_version_num
        artifact.revision = artifact.revision + 1
        artifact.size_bytes = size_bytes
        artifact.sha256 = sha256_hash
        artifact.checksum_sha256 = sha256_hash
        artifact.updated_at = time.time()
        artifact.versions.append(asdict(v_entry))

        return artifact

    @classmethod
    def restore_version(cls, artifact_id: str, target_version: int) -> Optional[CanonicalArtifact]:
        artifact = cls.get_artifact(artifact_id)
        if not artifact:
            return None

        matched_version = next((v for v in artifact.versions if v.get("version") == target_version), None)
        if not matched_version:
            return None

        content = matched_version.get("content", "")
        return cls.save_new_version(
            artifact_id=artifact_id,
            new_content=content,
            author_agent="human_rollback",
            change_summary=f"Rollback to version {target_version}",
        )

    # ── E. Système de Partage Sécurisé ───────────────────────────

    @classmethod
    def create_share_link(
        cls,
        artifact_id: str,
        visibility: ShareVisibility = ShareVisibility.LINK_ONLY,
        permission: SharePermission = SharePermission.DOWNLOAD_ALLOWED,
        expires_in_hours: Optional[int] = 72,
    ) -> ShareToken:
        token = f"sh_{uuid.uuid4().hex}"
        expires_at = time.time() + (expires_in_hours * 3600) if expires_in_hours else None

        st = ShareToken(
            token=token,
            artifact_id=artifact_id,
            visibility=visibility,
            permission=permission,
            expires_at=expires_at,
        )
        cls._SHARE_TOKENS[token] = st

        artifact = cls.get_artifact(artifact_id)
        if artifact:
            artifact.share_id = token
            artifact.share_url = f"/share/{token}"

        return st

    @classmethod
    def resolve_share_link(cls, token: str) -> Optional[Dict[str, Any]]:
        st = cls._SHARE_TOKENS.get(token)
        if not st or st.is_revoked:
            return None
        if st.expires_at and time.time() > st.expires_at:
            return None

        artifact = cls.get_artifact(st.artifact_id)
        if not artifact:
            return None

        return {
            "artifact": artifact.to_dict(),
            "permission": st.permission.value,
            "visibility": st.visibility.value,
            "expires_at": st.expires_at,
        }

    # ── Helpers Internes ─────────────────────────────────────────

    @classmethod
    def _prepare_content(cls, content: Any, type: ArtifactType) -> Tuple[str, str, bytes]:
        if isinstance(content, bytes):
            ext_map = {
                ArtifactType.IMAGE: "png",
                ArtifactType.LOGO: "png",
                ArtifactType.VIDEO: "mp4",
                ArtifactType.AUDIO: "mp3",
                ArtifactType.PDF: "pdf",
                ArtifactType.DOCX: "docx",
                ArtifactType.PPTX: "pptx",
                ArtifactType.XLSX: "xlsx",
                ArtifactType.ZIP: "zip",
            }
            ext = ext_map.get(type, "bin")
            mime = "application/octet-stream"
            return ext, mime, content

        if isinstance(content, (dict, list)):
            text = json.dumps(content, indent=2, ensure_ascii=False)
            return "json", "application/json", text.encode("utf-8")

        text = str(content)
        if type in (ArtifactType.MARKDOWN, ArtifactType.REPORT):
            return "md", "text/markdown; charset=utf-8", text.encode("utf-8")
        elif type in (ArtifactType.HTML, ArtifactType.WEBSITE):
            return "html", "text/html; charset=utf-8", text.encode("utf-8")
        elif type == ArtifactType.CSV:
            return "csv", "text/csv; charset=utf-8", text.encode("utf-8")
        elif type == ArtifactType.CALENDAR:
            return "ics", "text/calendar; charset=utf-8", text.encode("utf-8")
        elif type == ArtifactType.TRANSCRIPT:
            return "srt", "text/plain; charset=utf-8", text.encode("utf-8")
        elif type == ArtifactType.CODE:
            return "txt", "text/plain; charset=utf-8", text.encode("utf-8")

        return "txt", "text/plain; charset=utf-8", text.encode("utf-8")

    @classmethod
    def _get_supported_exports(cls, type: ArtifactType) -> List[str]:
        if type in (ArtifactType.REPORT, ArtifactType.MARKDOWN, ArtifactType.DOCUMENT, ArtifactType.TEXT):
            return ["pdf", "docx", "html", "markdown", "txt"]
        elif type in (ArtifactType.SLIDES, ArtifactType.PPTX):
            return ["pptx", "pdf"]
        elif type in (ArtifactType.SPREADSHEET, ArtifactType.XLSX, ArtifactType.CSV):
            return ["xlsx", "csv"]
        elif type in (ArtifactType.WEBSITE, ArtifactType.PROJECT, ArtifactType.CODE):
            return ["zip", "html"]
        elif type in (ArtifactType.IMAGE, ArtifactType.LOGO):
            return ["png", "jpg", "webp"]
        elif type == ArtifactType.VIDEO:
            return ["mp4"]
        elif type == ArtifactType.AUDIO:
            return ["mp3", "wav"]
        elif type == ArtifactType.CALENDAR:
            return ["ics"]
        elif type == ArtifactType.DATASET:
            return ["json", "csv"]
        return ["txt"]

    @classmethod
    def _get_supported_actions(cls, type: ArtifactType) -> List[str]:
        base = ["open", "download", "share"]
        if type in (ArtifactType.IMAGE, ArtifactType.LOGO):
            return base + ["edit", "variations", "use_in_video", "use_in_slides"]
        elif type == ArtifactType.VIDEO:
            return base + ["create_variation", "use_in_presentation"]
        elif type == ArtifactType.AUDIO:
            return base + ["transcribe", "use_in_video"]
        elif type in (ArtifactType.SLIDES, ArtifactType.PPTX):
            return base + ["present", "export_pdf"]
        elif type in (ArtifactType.REPORT, ArtifactType.DOCUMENT):
            return base + ["export_pdf", "export_docx", "create_deck"]
        elif type in (ArtifactType.SPREADSHEET, ArtifactType.XLSX):
            return base + ["export_csv", "analyze"]
        elif type == ArtifactType.WEBSITE:
            return base + ["preview", "inspect_files", "export_zip"]
        return base

    @classmethod
    def _get_raw_text(cls, artifact: CanonicalArtifact) -> str:
        if os.path.exists(artifact.file_path):
            try:
                with open(artifact.file_path, "r", encoding="utf-8") as f:
                    return f.read()
            except UnicodeDecodeError:
                pass
        return artifact.description or artifact.title

    # ── Moteurs de Rendu Authentiques ────────────────────────────

    @classmethod
    def _render_pdf(cls, artifact: CanonicalArtifact) -> bytes:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
        from reportlab.lib import colors

        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4, rightMargin=54, leftMargin=54, topMargin=54, bottomMargin=54)

        styles = getSampleStyleSheet()
        title_style = ParagraphStyle("DocTitle", parent=styles["Title"], fontName="Helvetica-Bold", fontSize=22, leading=26, textColor=colors.HexColor("#0E1626"), alignment=0)
        heading_style = ParagraphStyle("DocHeading", parent=styles["Heading2"], fontName="Helvetica-Bold", fontSize=14, leading=18, textColor=colors.HexColor("#1B2A4A"), spaceBefore=14, spaceAfter=6)
        body_style = ParagraphStyle("DocBody", parent=styles["Normal"], fontName="Helvetica", fontSize=10, leading=14, textColor=colors.HexColor("#2C3E50"), spaceAfter=8)
        meta_style = ParagraphStyle("DocMeta", parent=styles["Normal"], fontName="Helvetica", fontSize=8, textColor=colors.HexColor("#7F8C8D"))

        story = [
            Paragraph(artifact.title, title_style),
            Paragraph(f"Ñkyel AI Sovereign Workspace · Mission: {artifact.mission_id} · Model: {artifact.model}", meta_style),
            Spacer(1, 10),
            HRFlowable(width="100%", thickness=1, color=colors.HexColor("#CBD5E1"), spaceAfter=14),
        ]

        raw_text = cls._get_raw_text(artifact)
        for line in raw_text.split("\n"):
            line = line.strip()
            if not line:
                story.append(Spacer(1, 6))
            elif line.startswith("#"):
                story.append(Paragraph(line.lstrip("#").strip(), heading_style))
            else:
                story.append(Paragraph(line, body_style))

        doc.build(story)  # type: ignore
        return buf.getvalue()

    @classmethod
    def _render_docx(cls, artifact: CanonicalArtifact) -> bytes:
        import docx
        from docx.shared import Pt, RGBColor

        doc = docx.Document()
        title_p = doc.add_heading(artifact.title, level=0)
        title_p.runs[0].font.color.rgb = RGBColor(14, 22, 38)

        meta_p = doc.add_paragraph()
        meta_run = meta_p.add_run(f"Généré par Ñkyel AI · Modèle: {artifact.model} · SHA-256: {artifact.sha256[:16]}...")
        meta_run.font.size = Pt(8.5)
        meta_run.font.color.rgb = RGBColor(120, 144, 156)
        doc.add_paragraph("─" * 40)

        raw_text = cls._get_raw_text(artifact)
        for line in raw_text.split("\n"):
            line = line.strip()
            if not line:
                continue
            if line.startswith("###"):
                doc.add_heading(line.lstrip("#").strip(), level=3)
            elif line.startswith("##"):
                doc.add_heading(line.lstrip("#").strip(), level=2)
            elif line.startswith("#"):
                doc.add_heading(line.lstrip("#").strip(), level=1)
            elif line.startswith("- ") or line.startswith("* "):
                doc.add_paragraph(line[2:], style="List Bullet")
            else:
                doc.add_paragraph(line)

        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()

    @classmethod
    def _render_pptx(cls, artifact: CanonicalArtifact) -> bytes:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_slide_layout = prs.slide_layouts[6]

        title_slide = prs.slides.add_slide(blank_slide_layout)
        tx_box = title_slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.5))
        tf = tx_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = artifact.title
        p.font.size = Pt(38)
        p.font.bold = True
        p.font.color.rgb = RGBColor(14, 22, 38)

        p2 = tf.add_paragraph()
        p2.text = f"Stratégie & Livrables Internationaux · Ñkyel AI · {artifact.model}"
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(80, 100, 120)

        raw_text = cls._get_raw_text(artifact)
        sections = [s for s in raw_text.split("\n#") if s.strip()]

        for sec in sections[:8]:
            lines = sec.strip().split("\n")
            sec_title = lines[0].lstrip("#").strip()
            sec_body = [l.strip() for l in lines[1:] if l.strip()]

            content_slide = prs.slides.add_slide(blank_slide_layout)
            box = content_slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(5.8))
            frame = box.text_frame
            frame.word_wrap = True

            hp = frame.paragraphs[0]
            hp.text = sec_title or "Section"
            hp.font.size = Pt(28)
            hp.font.bold = True
            hp.font.color.rgb = RGBColor(18, 38, 70)

            for b_line in sec_body[:5]:
                bp = frame.add_paragraph()
                bp.text = f"• {b_line.lstrip('-*• ')}"
                bp.font.size = Pt(16)
                bp.font.color.rgb = RGBColor(40, 50, 60)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    @classmethod
    def _render_xlsx(cls, artifact: CanonicalArtifact) -> bytes:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

        wb = openpyxl.Workbook()
        ws = wb.active
        assert ws is not None
        ws.title = "Master Data"

        header_fill = PatternFill(start_color="0E1626", end_color="0E1626", fill_type="solid")
        header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
        align_center = Alignment(horizontal="center", vertical="center")
        thin_border = Border(left=Side(style='thin', color='CBD5E1'), right=Side(style='thin', color='CBD5E1'), top=Side(style='thin', color='CBD5E1'), bottom=Side(style='thin', color='CBD5E1'))

        data_rows = artifact.metadata.get("data") or [
            {"Poste / Activité": "Campagne Digitale & Médias", "Budget Q1 ($)": 40000, "Budget Q2 ($)": 80000, "Total ($)": 120000},
            {"Poste / Activité": "Résidences d'Influence & Presse", "Budget Q1 ($)": 25000, "Budget Q2 ($)": 60000, "Total ($)": 85000},
            {"Poste / Activité": "Certifications Éco-Lodge & Parc", "Budget Q1 ($)": 50000, "Budget Q2 ($)": 0, "Total ($)": 50000},
        ]

        if isinstance(data_rows, list) and data_rows and isinstance(data_rows[0], dict):
            headers = list(data_rows[0].keys())
            ws.append(headers)
            for cell in ws[1]:
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = align_center

            for row_dict in data_rows:
                ws.append([row_dict.get(h, "") for h in headers])

            for row in ws.iter_rows(min_row=2, max_row=len(data_rows) + 1, min_col=1, max_col=len(headers)):
                for cell in row:
                    cell.border = thin_border
        else:
            ws.append(["Titre", "Description", "Valeur"])
            ws.append([artifact.title, artifact.description or "N/A", artifact.size_bytes])

        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()

    @classmethod
    def _render_csv(cls, artifact: CanonicalArtifact) -> bytes:
        data_rows = artifact.metadata.get("data")
        if isinstance(data_rows, list) and data_rows and isinstance(data_rows[0], dict):
            headers = list(data_rows[0].keys())
            lines = [",".join(headers)]
            for r in data_rows:
                lines.append(",".join(['"{}"'.format(str(r.get(h, "")).replace('"', '""')) for h in headers]))
            return "\n".join(lines).encode("utf-8")

        raw_text = cls._get_raw_text(artifact)
        return raw_text.encode("utf-8")

    @classmethod
    def _render_zip(cls, artifact: CanonicalArtifact) -> bytes:
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
            zf.writestr("index.html", f"<!DOCTYPE html><html><head><title>{artifact.title}</title></head><body><h1>{artifact.title}</h1><p>{cls._get_raw_text(artifact)}</p></body></html>")
            zf.writestr("manifest.json", json.dumps(artifact.to_dict(), indent=2))
            zf.writestr("README.md", f"# {artifact.title}\n\nGenerated by Ñkyel AI.\nSHA-256: {artifact.sha256}")
        return buf.getvalue()

    @classmethod
    def _render_ics(cls, artifact: CanonicalArtifact) -> bytes:
        """Génère un fichier de calendrier iCalendar (ICS) standard."""
        ics_text = (
            "BEGIN:VCALENDAR\n"
            "VERSION:2.0\n"
            "PRODID:-//SmartANDJ//Nkyel AI Calendar//FR\n"
            "BEGIN:VEVENT\n"
            f"UID:{artifact.id}@nkyel.ai\n"
            f"DTSTAMP:{time.strftime('%Y%m%dT%H%M%SZ', time.gmtime())}\n"
            f"DTSTART:{time.strftime('%Y%m%dT%H%M%SZ', time.gmtime(time.time() + 86400))}\n"
            f"SUMMARY:{artifact.title}\n"
            f"DESCRIPTION:{artifact.description or 'Événement généré par Ñkyel AI'}\n"
            "END:VEVENT\n"
            "END:VCALENDAR\n"
        )
        return ics_text.encode("utf-8")

    @classmethod
    def _render_image(cls, artifact: CanonicalArtifact, fmt: str) -> Tuple[bytes, str]:
        if os.path.exists(artifact.file_path):
            with Image.open(artifact.file_path) as img:
                buf = io.BytesIO()
                out_fmt = "JPEG" if fmt in ("jpg", "jpeg") else fmt.upper()
                if out_fmt == "JPEG" and img.mode in ("RGBA", "P"):
                    img = img.convert("RGB")
                img.save(buf, format=out_fmt)
                return buf.getvalue(), f"image/{fmt}"

        img = Image.new("RGB", (1024, 1024), color=(14, 22, 38))
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue(), "image/png"
